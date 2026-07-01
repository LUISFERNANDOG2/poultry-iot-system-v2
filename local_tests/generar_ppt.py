from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de colores ──────────────────────────────────────────────────────────
VERDE_OSCURO  = RGBColor(0x1B, 0x5E, 0x20)   # fondo título
VERDE_MEDIO   = RGBColor(0x2E, 0x7D, 0x32)   # encabezados slide
VERDE_CLARO   = RGBColor(0x81, 0xC7, 0x84)   # acento
AMARILLO      = RGBColor(0xFF, 0xD6, 0x00)   # alerta amarilla
ROJO          = RGBColor(0xC6, 0x28, 0x28)   # alerta roja
AZUL_DATO     = RGBColor(0x15, 0x65, 0xC0)   # valores numéricos
GRIS_FONDO    = RGBColor(0xF4, 0xF6, 0xF8)   # fondo slides normales
BLANCO        = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO         = RGBColor(0x1A, 0x1A, 0x1A)
GRIS_TEXTO    = RGBColor(0x42, 0x42, 0x42)


def rgb(color):
    return color


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=14, bold=False, color=NEGRO,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_paragraph(tf, text, font_size=12, bold=False, color=NEGRO,
                  align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def title_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, VERDE_OSCURO)

    # Franja decorativa inferior
    add_rect(slide, 0, 6.5, 10, 1.0, fill_color=VERDE_MEDIO)

    # Ícono/logo placeholder — círculo verde claro
    add_rect(slide, 0.35, 0.3, 1.0, 1.0, fill_color=VERDE_CLARO)
    add_text_box(slide, "🐔", 0.35, 0.22, 1.0, 1.0,
                 font_size=32, color=BLANCO, align=PP_ALIGN.CENTER)

    # Título principal
    add_text_box(slide, "Sistema de Monitoreo IoT\nPara Caseta Avícola",
                 1.5, 0.25, 8.0, 1.4,
                 font_size=32, bold=True, color=BLANCO)

    # Subtítulo
    add_text_box(slide, "Análisis de datos reales\nMódulo M1 — 3 al 10 de marzo 2025",
                 1.5, 1.75, 8.0, 1.0,
                 font_size=18, color=VERDE_CLARO)

    # Línea separadora
    add_rect(slide, 1.5, 2.85, 7.0, 0.04, fill_color=VERDE_CLARO)

    # Descripción
    add_text_box(slide,
                 "309,781 lecturas de sensores de temperatura, humedad, CO₂, CO y NH₃\n"
                 "Valor de las alertas automáticas y detección de anomalías",
                 1.5, 3.05, 8.0, 1.2,
                 font_size=14, color=BLANCO)

    # Footer
    add_text_box(slide, "Proyecto de Avicultura de Precisión  |  Datos reales de campo",
                 0, 6.6, 10, 0.4,
                 font_size=11, color=VERDE_CLARO, align=PP_ALIGN.CENTER)
    return slide


def section_divider(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, VERDE_MEDIO)
    add_rect(slide, 0, 2.8, 10, 0.08, fill_color=AMARILLO)
    add_text_box(slide, title, 0.5, 2.0, 9.0, 1.2,
                 font_size=36, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, subtitle, 0.5, 3.2, 9.0, 0.8,
                     font_size=18, color=VERDE_CLARO, align=PP_ALIGN.CENTER)
    return slide


def resumen_datos(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, GRIS_FONDO)

    # Header
    add_rect(slide, 0, 0, 10, 1.0, fill_color=VERDE_MEDIO)
    add_text_box(slide, "Resumen del conjunto de datos",
                 0.3, 0.1, 9.4, 0.8, font_size=24, bold=True, color=BLANCO)

    # 4 tarjetas de resumen
    cards = [
        ("309,781",   "lecturas totales",   VERDE_MEDIO),
        ("1 módulo",  "M1 — caseta única",  AZUL_DATO),
        ("7 días",    "3–10 marzo 2025",    VERDE_OSCURO),
        ("5 variables", "Temp · Hum · CO₂ · CO · NH₃", GRIS_TEXTO),
    ]
    for i, (val, label, col) in enumerate(cards):
        x = 0.3 + i * 2.35
        add_rect(slide, x, 1.2, 2.1, 1.4, fill_color=col)
        add_text_box(slide, val, x, 1.25, 2.1, 0.75,
                     font_size=22, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x, 1.95, 2.1, 0.55,
                     font_size=11, color=BLANCO, align=PP_ALIGN.CENTER)

    # Nota de sensores de gas
    add_rect(slide, 0.3, 2.85, 9.4, 0.65, fill_color=RGBColor(0xFF, 0xF9, 0xC4), line_color=AMARILLO, line_width=Pt(1.5))
    add_text_box(slide,
                 "⚠  Nota sobre sensores de gas (CO, CO₂, NH₃): probablemente sin calibración. "
                 "Los valores absolutos son solo de referencia — los patrones y tendencias sí son válidos.",
                 0.4, 2.9, 9.2, 0.55, font_size=11, color=RGBColor(0x5D, 0x40, 0x00))

    # Tabla de variables
    headers = ["Variable", "Promedio semanal", "Mínimo", "Máximo", "Estado calibración"]
    rows = [
        ("Temperatura (°C)",  "31.98°C",    "28.0°C",    "35.1°C",    "✔ Confiable"),
        ("Humedad (%HR)",     "34.1%",      "21%",       "58%",       "✔ Confiable"),
        ("CO₂ (sensor)",      "2,018 u",    "400 u",     "20,823 u",  "⚠ Sin calibrar"),
        ("CO (sensor)",       "1,374 u",    "971 u",     "2,587 u",   "⚠ Sin calibrar"),
        ("NH₃ / Amoniaco",    "2,430 u",    "1,799 u",   "3,856 u",   "⚠ Sin calibrar"),
    ]
    col_w = [2.3, 1.8, 1.2, 1.2, 1.9]
    col_x = [0.3, 2.6, 4.4, 5.6, 6.8]
    row_h = 0.38

    # Encabezado tabla
    for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
        add_rect(slide, cx, 3.65, cw - 0.05, row_h, fill_color=VERDE_OSCURO)
        add_text_box(slide, hdr, cx + 0.05, 3.68, cw - 0.1, row_h - 0.05,
                     font_size=10, bold=True, color=BLANCO)

    for ri, row in enumerate(rows):
        bg = BLANCO if ri % 2 == 0 else RGBColor(0xEC, 0xF5, 0xEC)
        for ci, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
            add_rect(slide, cx, 4.03 + ri * row_h, cw - 0.05, row_h - 0.02,
                     fill_color=bg, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
            fc = ROJO if "⚠" in val else (VERDE_MEDIO if "✔" in val else NEGRO)
            add_text_box(slide, val, cx + 0.05, 4.05 + ri * row_h, cw - 0.1, row_h - 0.04,
                         font_size=10, color=fc)
    return slide


def slide_temperatura(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, GRIS_FONDO)

    add_rect(slide, 0, 0, 10, 1.0, fill_color=RGBColor(0xBF, 0x36, 0x0C))
    add_text_box(slide, "Temperatura  —  Módulo M1", 0.3, 0.1, 9.4, 0.8,
                 font_size=24, bold=True, color=BLANCO)

    # KPIs
    kpis = [
        ("31.98°C", "Promedio semanal", RGBColor(0x2E, 0x7D, 0x32)),
        ("35.1°C",  "Máximo absoluto\n4 mar, 8:33am", RGBColor(0xC6, 0x28, 0x28)),
        ("28.0°C",  "Mínimo absoluto", AZUL_DATO),
        ("32–33°C", "Rango ideal\nprimera semana", RGBColor(0x55, 0x55, 0x55)),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        x = 0.2 + i * 2.4
        add_rect(slide, x, 1.1, 2.2, 1.3, fill_color=col)
        add_text_box(slide, val, x, 1.12, 2.2, 0.65,
                     font_size=20, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, x, 1.72, 2.2, 0.6,
                     font_size=10, color=BLANCO, align=PP_ALIGN.CENTER)

    # Tabla diaria
    add_text_box(slide, "Temperatura por día (hora local UTC-6):", 0.3, 2.55, 5.5, 0.35,
                 font_size=12, bold=True, color=VERDE_OSCURO)
    dias = ["3 mar", "4 mar", "5 mar", "6 mar", "7 mar", "8 mar", "9 mar", "10 mar"]
    promedios = [32.38, 32.48, 31.86, 32.02, 31.70, 31.40, 31.42, 31.02]
    maximos   = [35.0,  35.1,  33.9,  33.5,  33.3,  33.3,  33.8,  33.0]
    minimos   = [28.0,  30.0,  29.0,  30.0,  29.0,  29.0,  28.0,  29.0]

    row_h = 0.32
    headers2 = ["Día", "Prom °C", "Máx °C", "Mín °C", "Estado"]
    col_w2 = [1.0, 0.9, 0.9, 0.9, 1.4]
    col_x2 = [0.25, 1.25, 2.15, 3.05, 3.95]

    for ci, (hdr, cw, cx) in enumerate(zip(headers2, col_w2, col_x2)):
        add_rect(slide, cx, 2.95, cw - 0.04, row_h, fill_color=RGBColor(0xBF, 0x36, 0x0C))
        add_text_box(slide, hdr, cx + 0.03, 2.97, cw - 0.06, row_h - 0.04,
                     font_size=9, bold=True, color=BLANCO)

    for ri, (d, pr, mx, mn) in enumerate(zip(dias, promedios, maximos, minimos)):
        y = 3.27 + ri * row_h
        bg = BLANCO if ri % 2 == 0 else RGBColor(0xFB, 0xE9, 0xE7)
        estado = "⚠ PICO CALOR" if mx >= 35.0 else ("✔ Normal" if pr <= 33.0 else "→ Límite")
        est_col = ROJO if "PICO" in estado else (VERDE_MEDIO if "Normal" in estado else AMARILLO)
        vals = [d, f"{pr}", f"{mx}", f"{mn}", estado]
        for ci, (v, cw, cx) in enumerate(zip(vals, col_w2, col_x2)):
            add_rect(slide, cx, y, cw - 0.04, row_h - 0.02,
                     fill_color=bg, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(0.5))
            fc = est_col if ci == 4 else (RGBColor(0xBF, 0x36, 0x0C) if v == f"{mx}" and mx >= 35.0 else NEGRO)
            add_text_box(slide, v, cx + 0.04, y + 0.02, cw - 0.08, row_h - 0.04, font_size=9, color=fc)

    # Panel lateral — sensor fallando
    add_rect(slide, 5.55, 2.55, 4.1, 4.85,
             fill_color=RGBColor(0xFF, 0xEB, 0xEE), line_color=ROJO, line_width=Pt(1.5))
    add_text_box(slide, "🔴  PROBLEMA DETECTADO:\nSensor de temperatura fallando",
                 5.65, 2.6, 3.9, 0.7, font_size=11, bold=True, color=ROJO)

    datos_null = [("3 mar", 4), ("4 mar", 19), ("5 mar", 46),
                  ("6 mar", 43), ("7 mar", 53), ("8 mar", 60), ("9 mar", 67)]
    add_text_box(slide, "Porcentaje de lecturas SIN temperatura:", 5.65, 3.35, 3.9, 0.3,
                 font_size=9, bold=True, color=NEGRO)

    bar_max_w = 3.0
    for ri, (d, pct) in enumerate(datos_null):
        y = 3.7 + ri * 0.41
        bar_w = bar_max_w * pct / 100
        add_text_box(slide, d, 5.65, y, 0.65, 0.35, font_size=9, color=NEGRO)
        # Fondo barra
        add_rect(slide, 6.3, y + 0.04, bar_max_w, 0.25, fill_color=RGBColor(0xEE, 0xEE, 0xEE))
        # Barra rellena
        bar_col = ROJO if pct >= 50 else (AMARILLO if pct >= 30 else VERDE_CLARO)
        add_rect(slide, 6.3, y + 0.04, bar_w, 0.25, fill_color=bar_col)
        add_text_box(slide, f"{pct}%", 6.3 + bar_w + 0.05, y, 0.5, 0.35, font_size=9,
                     bold=True, color=ROJO if pct >= 50 else NEGRO)

    add_text_box(slide, "→ Alerta \"sin datos\" hubiera detectado\n   la falla desde el día 5",
                 5.65, 7.0, 3.9, 0.45, font_size=9, italic=True, color=GRIS_TEXTO)
    return slide


def slide_humedad(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, GRIS_FONDO)

    add_rect(slide, 0, 0, 10, 1.0, fill_color=AZUL_DATO)
    add_text_box(slide, "Humedad Relativa  —  El hallazgo más crítico", 0.3, 0.1, 9.4, 0.8,
                 font_size=24, bold=True, color=BLANCO)

    kpis = [
        ("34.1%",  "Promedio semanal\n(Recomendado: 50–70%)", ROJO),
        ("21%",    "Mínimo absoluto\n7 mar · 9:54am", ROJO),
        ("58%",    "Máximo absoluto\n2 mar noche", VERDE_MEDIO),
        ("50–70%", "Rango óptimo\nprimera semana de vida", AZUL_DATO),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        x = 0.2 + i * 2.4
        add_rect(slide, x, 1.1, 2.2, 1.3, fill_color=col)
        add_text_box(slide, val, x, 1.12, 2.2, 0.65,
                     font_size=20, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, x, 1.72, 2.2, 0.6,
                     font_size=10, color=BLANCO, align=PP_ALIGN.CENTER)

    # Patrón horario
    add_text_box(slide, "Patrón horario de humedad (promedio semanal):", 0.3, 2.55, 9.4, 0.35,
                 font_size=12, bold=True, color=VERDE_OSCURO)

    # Gráfica de barras horaria simplificada
    horas_hum = {
        6: 31, 7: 29, 8: 27, 9: 26, 10: 26, 11: 27,
        12: 29, 13: 32, 14: 34, 15: 35, 16: 36, 17: 37,
        18: 37, 19: 36, 20: 36, 21: 37, 22: 38, 23: 39,
    }
    bar_area_x = 0.3
    bar_area_y = 3.0
    bar_area_w = 9.4
    bar_area_h = 2.3
    n = len(horas_hum)
    bw = bar_area_w / n - 0.04

    add_text_box(slide, "% HR", 0.3, 2.9, 0.5, 0.25, font_size=8, color=GRIS_TEXTO)
    # línea de 30% (alerta)
    y_30 = bar_area_y + bar_area_h - (bar_area_h * 30 / 60)
    add_rect(slide, bar_area_x, y_30, bar_area_w, 0.02, fill_color=AMARILLO)
    add_text_box(slide, "⚠ 30% alerta", bar_area_x + 7.0, y_30 - 0.2, 1.5, 0.25,
                 font_size=8, color=AMARILLO)
    # línea de 25% (crítico)
    y_25 = bar_area_y + bar_area_h - (bar_area_h * 25 / 60)
    add_rect(slide, bar_area_x, y_25, bar_area_w, 0.02, fill_color=ROJO)
    add_text_box(slide, "🔴 25% crítico", bar_area_x + 7.0, y_25 - 0.2, 1.5, 0.25,
                 font_size=8, color=ROJO)

    for i, (hora, val) in enumerate(horas_hum.items()):
        bx = bar_area_x + i * (bw + 0.04)
        bh = bar_area_h * val / 60
        by = bar_area_y + bar_area_h - bh
        col = ROJO if val < 25 else (AMARILLO if val < 30 else AZUL_DATO)
        add_rect(slide, bx, by, bw, bh, fill_color=col)
        if i % 3 == 0:
            add_text_box(slide, f"{hora}h", bx, bar_area_y + bar_area_h + 0.02, bw + 0.1, 0.25,
                         font_size=7, color=GRIS_TEXTO, align=PP_ALIGN.CENTER)

    # Nota explicativa
    add_rect(slide, 0.3, 5.5, 9.4, 1.4,
             fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=AZUL_DATO, line_width=Pt(1))
    add_text_box(slide,
                 "Patrón observado:  La humedad cae fuertemente entre las 7–10am al subir la temperatura "
                 "(sol matutino + calefactores). Alcanza su mínimo justo cuando los pollitos están más activos.\n\n"
                 "Consecuencias de humedad < 25%:  deshidratación, irritación de mucosas respiratorias, mayor susceptibilidad "
                 "a enfermedades respiratorias y menor conversión alimenticia.\n\n"
                 "Corrección posible:  nebulizadores automáticos activados por alerta + revisión de ventilación diurna.",
                 0.45, 5.55, 9.1, 1.3, font_size=10, color=NEGRO)
    return slide


def slide_co2(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, GRIS_FONDO)

    add_rect(slide, 0, 0, 10, 1.0, fill_color=RGBColor(0x4A, 0x14, 0x8C))
    add_text_box(slide, "CO₂  —  Patrón de ventilación y acumulación", 0.3, 0.1, 9.4, 0.8,
                 font_size=24, bold=True, color=BLANCO)

    add_rect(slide, 0.3, 1.05, 9.4, 0.45,
             fill_color=RGBColor(0xFF, 0xF9, 0xC4), line_color=AMARILLO, line_width=Pt(1))
    add_text_box(slide,
                 "⚠  Sensor probablemente sin calibración — valores absolutos son de referencia. "
                 "Los patrones y tendencias son válidos.",
                 0.4, 1.08, 9.2, 0.38, font_size=10, color=RGBColor(0x5D, 0x40, 0x00))

    # Tabla diaria
    add_text_box(slide, "CO₂ por día — tendencia semanal:", 0.3, 1.65, 5.0, 0.35,
                 font_size=12, bold=True, color=RGBColor(0x4A, 0x14, 0x8C))
    dias_co2 = ["3 mar", "4 mar", "5 mar", "6 mar", "7 mar", "8 mar", "9 mar", "10 mar"]
    promedios_co2 = [1780, 1817, 1509, 1827, 2097, 2363, 2456, 3092]
    maximos_co2   = [12472, 4750, 2587, 3613, 4204, 4503, 4536, 4120]

    headers3 = ["Día", "CO₂ prom", "CO₂ máx", "Tendencia"]
    col_w3 = [1.1, 1.2, 1.2, 1.7]
    col_x3 = [0.25, 1.35, 2.55, 3.75]
    row_h = 0.33

    for ci, (hdr, cw, cx) in enumerate(zip(headers3, col_w3, col_x3)):
        add_rect(slide, cx, 2.05, cw - 0.04, row_h, fill_color=RGBColor(0x4A, 0x14, 0x8C))
        add_text_box(slide, hdr, cx + 0.03, 2.07, cw - 0.06, row_h - 0.04,
                     font_size=9, bold=True, color=BLANCO)

    for ri, (d, pr, mx) in enumerate(zip(dias_co2, promedios_co2, maximos_co2)):
        y = 2.38 + ri * row_h
        bg = BLANCO if ri % 2 == 0 else RGBColor(0xEE, 0xE8, 0xF7)
        if pr > 2500:
            tend = "↑↑ SUBIENDO"
            tc = ROJO
        elif pr > 2000:
            tend = "↑ Elevado"
            tc = RGBColor(0xE6, 0x51, 0x00)
        else:
            tend = "→ Estable"
            tc = VERDE_MEDIO
        vals = [d, f"{pr:,}", f"{mx:,}", tend]
        for ci, (v, cw, cx) in enumerate(zip(vals, col_w3, col_x3)):
            add_rect(slide, cx, y, cw - 0.04, row_h - 0.02,
                     fill_color=bg, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(0.5))
            fc = tc if ci == 3 else (ROJO if pr > 2500 and ci == 1 else NEGRO)
            add_text_box(slide, v, cx + 0.04, y + 0.02, cw - 0.08, row_h - 0.04, font_size=9, color=fc)

    # Panel evento crítico
    add_rect(slide, 5.55, 1.65, 4.1, 3.3,
             fill_color=RGBColor(0xF3, 0xE5, 0xF5), line_color=RGBColor(0x4A, 0x14, 0x8C), line_width=Pt(1.5))
    add_text_box(slide, "🔴  EVENTO CRÍTICO\n2–3 marzo, ~23:00h",
                 5.65, 1.7, 3.9, 0.6, font_size=12, bold=True, color=ROJO)
    add_text_box(slide,
                 "CO₂ llegó a 20,823 unidades — el\nmáximo absoluto de toda la semana.\n\n"
                 "Primera noche de operación con:\n"
                 "  • Calefactores encendidos\n"
                 "  • Caseta completamente cerrada\n"
                 "  • Sin ventilación activa\n\n"
                 "CO y NH₃ también alcanzaron\nsus máximos esa misma noche.",
                 5.65, 2.4, 3.9, 2.4, font_size=10, color=NEGRO)

    # Patrón noche vs día
    add_rect(slide, 0.25, 5.08, 5.3, 1.3,
             fill_color=RGBColor(0xED, 0xE7, 0xF6), line_color=RGBColor(0x4A, 0x14, 0x8C), line_width=Pt(1))
    add_text_box(slide, "Patrón circadiano de CO₂:", 0.35, 5.12, 5.1, 0.3,
                 font_size=11, bold=True, color=RGBColor(0x4A, 0x14, 0x8C))
    add_text_box(slide,
                 "Noche (19–5h):  2,265 promedio  |  Caseta cerrada, sin ventilación\n"
                 "Día (6–18h):    1,798 promedio  |  Ventilación activa, CO₂ baja\n\n"
                 "→ Picos más bajos: 7–9am (725 u) cuando se abren cortinas\n"
                 "→ Picos más altos: 17–18h y madrugada (>2,700 u)",
                 0.35, 5.45, 5.1, 0.85, font_size=9, color=NEGRO)

    add_rect(slide, 5.55, 5.08, 4.1, 1.3,
             fill_color=RGBColor(0xFF, 0xEB, 0xEE), line_color=ROJO, line_width=Pt(1))
    add_text_box(slide, "Tendencia al alza días 7–10:", 5.65, 5.12, 3.9, 0.3,
                 font_size=11, bold=True, color=ROJO)
    add_text_box(slide,
                 "La cama húmeda y acumulación de excreta\nelevaron el CO₂ promedio diario de\n"
                 "1,500–1,800 a 3,092 unidades al final.\n\n"
                 "→ Alerta de tendencia hubiera notificado\n   1–2 días antes.",
                 5.65, 5.45, 3.9, 0.85, font_size=9, color=NEGRO)
    return slide


def slide_gases_ref(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, GRIS_FONDO)

    add_rect(slide, 0, 0, 10, 1.0, fill_color=RGBColor(0x4E, 0x34, 0x2E))
    add_text_box(slide, "CO y NH₃  —  Solo referencia de patrones", 0.3, 0.1, 9.4, 0.8,
                 font_size=24, bold=True, color=BLANCO)

    add_rect(slide, 0.3, 1.05, 9.4, 0.55,
             fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=RGBColor(0xE6, 0x51, 0x00), line_width=Pt(1.5))
    add_text_box(slide,
                 "⚠  Los sensores de CO y NH₃ estaban probablemente sin calibración. "
                 "Los valores medidos no corresponden a ppm reales.\n"
                 "     Se presentan los patrones y tendencias únicamente — los valores absolutos no deben tomarse como referencia de seguridad.",
                 0.45, 1.08, 9.1, 0.48, font_size=10, color=RGBColor(0x5D, 0x40, 0x00))

    # Dos columnas
    # CO
    add_rect(slide, 0.3, 1.75, 4.55, 5.2,
             fill_color=BLANCO, line_color=RGBColor(0x4E, 0x34, 0x2E), line_width=Pt(1))
    add_rect(slide, 0.3, 1.75, 4.55, 0.45, fill_color=RGBColor(0x4E, 0x34, 0x2E))
    add_text_box(slide, "Monóxido de Carbono (CO)", 0.4, 1.78, 4.35, 0.38,
                 font_size=13, bold=True, color=BLANCO)

    co_stats = [
        ("Promedio semanal", "1,374 u"),
        ("Mínimo", "971 u"),
        ("Máximo", "2,587 u  (2 mar noche)"),
        ("Std. desv.", "165 u"),
        ("Límite real calibrado*", "< 10 ppm (humanos)"),
    ]
    for ri, (k, v) in enumerate(co_stats):
        add_text_box(slide, k + ":", 0.45, 2.3 + ri * 0.42, 2.3, 0.38,
                     font_size=10, bold=True, color=GRIS_TEXTO)
        add_text_box(slide, v, 2.75, 2.3 + ri * 0.42, 2.0, 0.38,
                     font_size=10, color=NEGRO)

    add_text_box(slide,
                 "Patrón observado:\n"
                 "  • Más alto en la noche (>1,500 u)\n"
                 "  • Más bajo a media mañana (~1,100 u)\n"
                 "  • El día 2 mar fue el más alto de la semana\n"
                 "  • Sigue el mismo ritmo que CO₂\n\n"
                 "Con calibración correcta, una alerta de\nCO sería la más urgente — es tóxico.",
                 0.45, 4.2, 4.3, 2.6, font_size=9.5, color=NEGRO)

    # NH3
    add_rect(slide, 5.15, 1.75, 4.55, 5.2,
             fill_color=BLANCO, line_color=RGBColor(0x1B, 0x5E, 0x20), line_width=Pt(1))
    add_rect(slide, 5.15, 1.75, 4.55, 0.45, fill_color=RGBColor(0x1B, 0x5E, 0x20))
    add_text_box(slide, "Amoniaco / NH₃", 5.25, 1.78, 4.35, 0.38,
                 font_size=13, bold=True, color=BLANCO)

    nh3_stats = [
        ("Promedio semanal", "2,430 u"),
        ("Mínimo", "1,799 u"),
        ("Máximo", "3,856 u  (2 mar noche)"),
        ("Std. desv.", "243 u"),
        ("Límite real calibrado*", "< 25 ppm para aves"),
    ]
    for ri, (k, v) in enumerate(nh3_stats):
        add_text_box(slide, k + ":", 5.3, 2.3 + ri * 0.42, 2.3, 0.38,
                     font_size=10, bold=True, color=GRIS_TEXTO)
        add_text_box(slide, v, 7.6, 2.3 + ri * 0.42, 2.0, 0.38,
                     font_size=10, color=NEGRO)

    add_text_box(slide,
                 "Patrón observado:\n"
                 "  • Relativamente estable toda la semana\n"
                 "  • Pico: 3,856 u la noche del 2–3 mar\n"
                 "  • Ligero aumento días 8–10 (cama húmeda)\n"
                 "  • Mínimos entre 7–9am, igual que CO₂\n\n"
                 "El NH₃ real en caseta avícola sube con la\n"
                 "cama húmeda y falta de ventilación — el\n"
                 "patrón detectado es coherente con eso.",
                 5.3, 4.2, 4.3, 2.6, font_size=9.5, color=NEGRO)

    add_text_box(slide,
                 "* Con sensores calibrados, el sistema puede emitir alertas reales de seguridad para CO y NH₃",
                 0.3, 7.05, 9.4, 0.35, font_size=9, italic=True, color=GRIS_TEXTO)
    return slide


def slide_eventos_criticos(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, GRIS_FONDO)

    add_rect(slide, 0, 0, 10, 1.0, fill_color=ROJO)
    add_text_box(slide, "Eventos críticos detectados — qué hubiera alertado el sistema",
                 0.3, 0.1, 9.4, 0.8, font_size=22, bold=True, color=BLANCO)

    eventos = [
        {
            "num": "1",
            "titulo": "Primera noche — caseta sin ventilación",
            "fecha": "2–3 mar, ~23:00–23:20h",
            "detalle": "CO₂: 20,823 u  |  CO: 2,587 u  |  NH₃: 3,856 u  |  Temp: 28°C",
            "causa": "Primera noche de operación: calefactores encendidos, caseta completamente cerrada.",
            "accion": "Alerta inmediata → abrir cortinas/ventilación + revisar calefactores",
            "col": ROJO,
        },
        {
            "num": "2",
            "titulo": "Pico de calor matutino",
            "fecha": "4 mar, 8:33–8:35am",
            "detalle": "Temperatura: 35.1°C  |  Humedad: 26%",
            "causa": "Sol matutino + calefactores aún activos = sobrecalentamiento.",
            "accion": "Alerta temp > 34°C → apagar calefactor + abrir ventilación",
            "col": RGBColor(0xBF, 0x36, 0x0C),
        },
        {
            "num": "3",
            "titulo": "Humedad crítica",
            "fecha": "4 mar 11:32am  y  7 mar 9:54am",
            "detalle": "Humedad: 21%  —  mínimo de la semana",
            "causa": "Temperatura alta + ventilación excesiva desecó el ambiente.",
            "accion": "Alerta humedad < 25% → activar nebulizadores",
            "col": AZUL_DATO,
        },
        {
            "num": "4",
            "titulo": "Sensor de temperatura fallando",
            "fecha": "A partir del 5 mar (progresivo)",
            "detalle": "Día 9 mar: 67% de lecturas sin dato de temperatura",
            "causa": "Sensor o conexión con falla intermitente, empeorando cada día.",
            "accion": "Alerta \"sin datos por X minutos\" → revisar hardware",
            "col": RGBColor(0x55, 0x55, 0x55),
        },
        {
            "num": "5",
            "titulo": "Acumulación de gases — fin de semana",
            "fecha": "8–10 mar (tendencia creciente)",
            "detalle": "CO₂ promedio: 2,363 → 2,456 → 3,092 en días consecutivos",
            "causa": "Acumulación de cama y excretas + ventilación insuficiente.",
            "accion": "Alerta de tendencia → renovar cama + aumentar ventilación",
            "col": RGBColor(0x4A, 0x14, 0x8C),
        },
    ]

    for i, ev in enumerate(eventos):
        y = 1.1 + i * 1.18
        add_rect(slide, 0.2, y, 0.5, 1.0, fill_color=ev["col"])
        add_text_box(slide, ev["num"], 0.2, y + 0.2, 0.5, 0.6,
                     font_size=20, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_rect(slide, 0.7, y, 9.1, 1.0,
                 fill_color=BLANCO, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(0.5))
        add_text_box(slide, ev["titulo"] + "  ·  " + ev["fecha"],
                     0.8, y + 0.03, 8.9, 0.3, font_size=11, bold=True, color=ev["col"])
        add_text_box(slide, ev["detalle"],
                     0.8, y + 0.32, 8.9, 0.25, font_size=10, color=NEGRO)
        add_text_box(slide, "Causa: " + ev["causa"] + "   →   " + ev["accion"],
                     0.8, y + 0.58, 8.9, 0.35, font_size=9, italic=True, color=GRIS_TEXTO)
    return slide


def slide_patron_circadiano(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, GRIS_FONDO)

    add_rect(slide, 0, 0, 10, 1.0, fill_color=VERDE_MEDIO)
    add_text_box(slide, "Patrón circadiano — comportamiento a lo largo del día",
                 0.3, 0.1, 9.4, 0.8, font_size=24, bold=True, color=BLANCO)

    # Tabla de patrones por franja horaria
    franjas = [
        ("Madrugada\n0–5h",
         "31.6°C",   "↑ 39%",  "↑ 2,200 u",  "↑ 2,550 u",
         "Caseta cerrada. Gases y humedad acumulados. CO₂ en máximo."),
        ("Amanecer\n6–9h",
         "↑ 32.1°C", "↓ 27%",  "↓ 750 u",    "↓ 2,060 u",
         "Se abren cortinas. Ventilación baja gases. Humedad cae con calor."),
        ("Media mañana\n9–12h",
         "PICO 33°C","MÍN 26%","→ 900 u",    "→ 2,100 u",
         "Temperatura máxima del día. Humedad mínima. Riesgo de estrés térmico."),
        ("Tarde\n13–18h",
         "↓ 31.5°C", "↑ 34%",  "↑ 2,600 u",  "↑ 2,580 u",
         "Temperatura baja. Humedad sube. Gases aumentan sin ventilación."),
        ("Noche\n19–23h",
         "31.6°C",   "↑ 37%",  "↑ 2,400 u",  "↑ 2,600 u",
         "Caseta se cierra. Gases y humedad escalan. Necesita ventilación."),
    ]

    headers4 = ["Franja", "Temp", "Humedad", "CO₂", "NH₃", "Observación"]
    col_w4 = [1.05, 0.75, 0.85, 0.85, 0.85, 4.95]
    col_x4 = [0.2, 1.25, 2.0, 2.85, 3.7, 4.55]
    row_h = 0.5

    for ci, (hdr, cw, cx) in enumerate(zip(headers4, col_w4, col_x4)):
        add_rect(slide, cx, 1.1, cw - 0.04, 0.38, fill_color=VERDE_OSCURO)
        add_text_box(slide, hdr, cx + 0.03, 1.12, cw - 0.06, 0.32,
                     font_size=10, bold=True, color=BLANCO)

    fila_cols = [
        RGBColor(0x1A, 0x23, 0x7E),  # madrugada azul oscuro
        RGBColor(0xF5, 0x7F, 0x17),  # amanecer naranja
        RGBColor(0xBF, 0x36, 0x0C),  # media mañana rojo
        RGBColor(0x1B, 0x5E, 0x20),  # tarde verde
        RGBColor(0x31, 0x1B, 0x92),  # noche violeta
    ]

    for ri, (franja, temp, hum, co2, nh3, obs) in enumerate(franjas):
        y = 1.48 + ri * (row_h + 0.06)
        bg = BLANCO if ri % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF0)
        vals = [franja, temp, hum, co2, nh3, obs]
        for ci, (v, cw, cx) in enumerate(zip(vals, col_w4, col_x4)):
            add_rect(slide, cx, y, cw - 0.04, row_h,
                     fill_color=fila_cols[ri] if ci == 0 else bg,
                     line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
            fc = BLANCO if ci == 0 else NEGRO
            add_text_box(slide, v, cx + 0.04, y + 0.04, cw - 0.08, row_h - 0.06,
                         font_size=9, color=fc, bold=(ci == 0))

    add_rect(slide, 0.2, 5.35, 9.6, 0.65,
             fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=VERDE_MEDIO, line_width=Pt(1))
    add_text_box(slide,
                 "Conclusión:  Las horas de mayor riesgo son 9–12h (calor + humedad mínima) y 0–5h (gases acumulados). "
                 "El sistema de alertas debería estar especialmente activo en esas franjas y al inicio de la noche "
                 "cuando se cierran las cortinas.",
                 0.35, 5.38, 9.3, 0.6, font_size=10, color=VERDE_OSCURO)
    return slide


def slide_alertas_propuestas(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, GRIS_FONDO)

    add_rect(slide, 0, 0, 10, 1.0, fill_color=VERDE_OSCURO)
    add_text_box(slide, "Alertas automáticas propuestas — basadas en los datos reales",
                 0.3, 0.1, 9.4, 0.8, font_size=22, bold=True, color=BLANCO)

    alertas = [
        ("🌡", "Temperatura alta",       "> 34°C por 3 min",       "Naranja", "Apagar calefactores, abrir ventilación"),
        ("🌡", "Temperatura crítica",     "> 35°C",                 "Rojo",    "Acción inmediata — riesgo de muerte por calor"),
        ("💧", "Humedad muy baja",        "< 30% por 10 min",       "Naranja", "Activar nebulizadores"),
        ("💧", "Humedad crítica",         "< 25%",                  "Rojo",    "Nebulizadores + revisar ventilación"),
        ("☁",  "CO₂ elevado",            "> 3,000 u por 5 min",    "Naranja", "Aumentar ventilación"),
        ("☁",  "CO₂ crítico",            "> 5,000 u",              "Rojo",    "Ventilación máxima + revisar calefactores"),
        ("📈", "Tendencia CO₂ al alza",   "Promedio horario +20%",  "Naranja", "Revisar cama + aumentar ventilación"),
        ("🔌", "Sensor sin datos",        "> 10 min sin lectura",   "Naranja", "Revisar hardware / conexión del módulo"),
    ]

    headers5 = ["", "Alerta", "Condición", "Nivel", "Acción recomendada"]
    col_w5 = [0.4, 1.9, 1.9, 0.8, 4.5]
    col_x5 = [0.2, 0.6, 2.5, 4.4, 5.2]
    row_h = 0.62

    for ci, (hdr, cw, cx) in enumerate(zip(headers5, col_w5, col_x5)):
        add_rect(slide, cx, 1.1, cw - 0.04, 0.35, fill_color=VERDE_OSCURO)
        add_text_box(slide, hdr, cx + 0.03, 1.12, cw - 0.06, 0.28,
                     font_size=10, bold=True, color=BLANCO)

    for ri, (icono, nombre, cond, nivel, accion) in enumerate(alertas):
        y = 1.45 + ri * (row_h + 0.03)
        bg = BLANCO if ri % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF0)
        nivel_col = ROJO if nivel == "Rojo" else RGBColor(0xE6, 0x51, 0x00)
        nivel_bg = RGBColor(0xFF, 0xEB, 0xEE) if nivel == "Rojo" else RGBColor(0xFF, 0xF3, 0xE0)
        vals = [icono, nombre, cond, nivel, accion]
        for ci, (v, cw, cx) in enumerate(zip(vals, col_w5, col_x5)):
            cell_bg = nivel_bg if ci == 3 else bg
            add_rect(slide, cx, y, cw - 0.04, row_h,
                     fill_color=cell_bg,
                     line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
            fc = nivel_col if ci == 3 else NEGRO
            fsz = 16 if ci == 0 else (10 if ci != 3 else 10)
            add_text_box(slide, v, cx + 0.04, y + 0.1, cw - 0.08, row_h - 0.12,
                         font_size=fsz, color=fc, bold=(ci == 3), align=(PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT))
    return slide


def slide_conclusiones(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, VERDE_OSCURO)

    add_text_box(slide, "Conclusiones y valor del sistema de alertas",
                 0.5, 0.2, 9.0, 0.8, font_size=26, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.0, 1.1, 8.0, 0.06, fill_color=VERDE_CLARO)

    puntos = [
        ("✔", "La temperatura fue manejada correctamente en general (31–33°C promedio), "
               "pero con un pico de 35.1°C el 4 de marzo que hubiera disparado alerta inmediata.",
         VERDE_CLARO),
        ("⚠", "La humedad estuvo consistentemente baja (34% promedio, 21% mínimo) — "
               "por debajo del rango recomendado para la primera semana. "
               "Las alertas automáticas de humedad son el beneficio más directo.",
         AMARILLO),
        ("⚠", "El sensor de temperatura empezó a fallar desde el día 5 y llegó a perder "
               "el 67% de las lecturas el día 9. Sin alerta de 'sensor sin datos', "
               "esto no se detecta a tiempo.",
         AMARILLO),
        ("🔴", "La primera noche fue la más crítica: todos los gases en sus máximos históricos "
                "con la caseta completamente cerrada. Una alerta ese momento hubiera evitado "
                "horas de exposición.",
         RGBColor(0xFF, 0x80, 0x80)),
        ("📈", "Los gases de CO₂ mostraron una tendencia creciente hacia el final de la semana "
                "(acumulación de cama). La detección de tendencias permite actuar antes de llegar "
                "al umbral crítico.",
         VERDE_CLARO),
        ("🔬", "Los sensores de CO y NH₃ requieren calibración para ser accionables. "
                "Con calibración correcta, el sistema daría alertas de seguridad completas.",
         RGBColor(0xAA, 0xD4, 0xFF)),
    ]

    for i, (icono, texto, color) in enumerate(puntos):
        y = 1.3 + i * 0.94
        add_rect(slide, 0.3, y, 0.55, 0.75, fill_color=color)
        add_text_box(slide, icono, 0.3, y + 0.1, 0.55, 0.55,
                     font_size=18, align=PP_ALIGN.CENTER, color=NEGRO)
        add_text_box(slide, texto, 0.95, y + 0.05, 8.7, 0.75,
                     font_size=11, color=BLANCO)
    return slide


def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    resumen_datos(prs)
    section_divider(prs, "Temperatura", "Bien controlada en general — con un evento de sobrecalentamiento y sensor fallando")
    slide_temperatura(prs)
    section_divider(prs, "Humedad", "El parámetro más crítico y accionable de la semana")
    slide_humedad(prs)
    section_divider(prs, "Calidad del Aire (CO₂, CO, NH₃)", "Patrones reales — valores absolutos de referencia por falta de calibración")
    slide_co2(prs)
    slide_gases_ref(prs)
    section_divider(prs, "Eventos Críticos", "Momentos donde las alertas automáticas hubieran hecho la diferencia")
    slide_eventos_criticos(prs)
    slide_patron_circadiano(prs)
    section_divider(prs, "Sistema de Alertas", "Propuesta basada en los datos reales de la semana")
    slide_alertas_propuestas(prs)
    slide_conclusiones(prs)

    out = r"D:\Poultry system\poultry-iot-system-v2\local_tests\analisis_caseta_M1_mar2025.pptx"
    prs.save(out)
    print(f"Guardado: {out}")


if __name__ == "__main__":
    main()
